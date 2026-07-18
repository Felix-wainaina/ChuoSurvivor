import base64
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from io import BytesIO
from PIL import Image
from pypdf import PdfReader
import docx
from pptx import Presentation
import io
import os
import json
from google import genai
from google.genai import types


app = FastAPI(title="Study Model API", description="API for processing images and text data.")


def load_env_file(env_path: str = ".env") -> None:
    """Load key/value pairs from a local .env file into os.environ."""
    if not os.path.exists(env_path):
        return

    with open(env_path, "r", encoding="utf-8") as env_file:
        for raw_line in env_file:
            line = raw_line.strip()
            if not line or line.startswith("#") or "=" not in line:
                continue

            key, value = line.split("=", 1)
            key = key.strip()
            value = value.strip().strip('"').strip("'")
            if key and key not in os.environ:
                os.environ[key] = value


load_env_file()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Google AI Studio Gemma 4 configuration.
MODEL_NAME = os.getenv("GEMINI_MODEL", "gemma-4-26b-a4b-it")
gemini_client: genai.Client | None = None


class QuizRequest(BaseModel):
    text: str


class StudyPlanRequest(BaseModel):
    messages: list[dict]


class ChatRequest(BaseModel):
    prompt: str


class Response(BaseModel):
    text: str


def get_gemini_client() -> genai.Client:
    global gemini_client
    if gemini_client is not None:
        return gemini_client

    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        raise HTTPException(
            status_code=500,
            detail="GEMINI_API_KEY is missing. Add it to backend/.env and restart the backend.",
        )
    # Keep one client alive for the FastAPI process. Creating it as a temporary
    # object can let its underlying HTTP client close during an async request.
    gemini_client = genai.Client(api_key=api_key)
    return gemini_client


@app.post("/chat", response_model=Response)
async def chat(payload: ChatRequest):
    """Cloud-first study chat used by the browser's offline fallback router."""
    if not payload.prompt.strip():
        raise HTTPException(status_code=400, detail="A study question is required.")
    try:
        response = await get_gemini_client().aio.models.generate_content(
            model=MODEL_NAME,
            contents=[types.Content(role="user", parts=[types.Part.from_text(text=payload.prompt)])],
            config=types.GenerateContentConfig(
                system_instruction=[types.Part.from_text(text=(
                    "You are a concise, encouraging study assistant. Answer using the supplied "
                    "study context, explain concepts clearly, and say when the context is insufficient."
                ))],
            ),
        )
        if not response.text:
            raise HTTPException(status_code=502, detail="Gemma returned an empty chat response.")
        return {"text": response.text.strip()}
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=502, detail=f"Gemma chat generation failed: {exc}")

def optimize_image_as_base64(file_bytes: bytes) -> str:
    """
    Optimize the image for better processing.
    This function can include resizing, converting to grayscale, etc.
    """
    try:
        img = Image.open(BytesIO(file_bytes))
        if img.mode != 'RGB':
            img = img.convert('RGB')
        # Resize the image to a maximum of 800x800 pixels while maintaining aspect ratio
        img.thumbnail((800, 800))
        buffered = BytesIO()
        img.save(buffered, format="JPEG", quality=85)  # Save as JPEG
        optimized_bytes = buffered.getvalue()
        return base64.b64encode(optimized_bytes).decode('utf-8')
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Invalid image file format: {str(e)}")
    
def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extracts raw text from all pages of a PDF."""
    try:
        pdf_file = io.BytesIO(file_bytes)
        reader = PdfReader(pdf_file)
        text_content = []
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_content.append(page_text)
        return "\n".join(text_content).strip()
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to parse PDF document: {str(e)}")


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extracts raw text from paragraphs in a .docx file."""
    try:
        docx_file = io.BytesIO(file_bytes)
        doc = docx.Document(docx_file)
        text_content = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(text_content).strip()
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to parse Word document: {str(e)}")


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extracts raw text from text shapes across all PowerPoint slides."""
    try:
        pptx_file = io.BytesIO(file_bytes)
        presentation = Presentation(pptx_file)
        text_content = []
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_content.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
        return "\n\n".join(text_content).strip()
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to parse PowerPoint presentation: {str(e)}")


# --- THE SUPERCHARGED ENDPOINT ---
@app.post("/explain", response_model=Response)
async def explain_content(
    file: UploadFile = File(...),
    user_prompt: str = Form(""),       # Optional custom focus prompt
    lang: str = Form("en")             # 'en' or 'sw'
):
    """Explain an uploaded file with Gemma 4 through Google AI Studio."""

    file_bytes = await file.read()
    if not file_bytes:
        raise HTTPException(status_code=400, detail="No file uploaded.")

    content_type = file.content_type or ""
    filename = (file.filename or "").lower()

    # 1. Map localization prompts
    if lang.lower() == "sw":
        system_instruction = (
            "Wewe ni mwalimu msaidizi anayesaidia kurahisisha dhana, picha, na maandishi. "
            "Eleza kile kilichowasilishwa kwa lugha rahisi sana ya Kiswahili, "
            "ukilenga kumsaidia mtu anayejifunza kuelewa kwa haraka."
        )
        fallback_prompt = "Tafadhali eleza maudhui haya kwa Kiswahili rahisi."
    else:
        system_instruction = (
            "You are an assistant teacher specializing in making complex charts, diagrams, or notes simple. "
            "Break down the provided material in highly simplified, plain English designed for absolute beginners."
        )
        fallback_prompt = "Please explain this content in simple English."

    final_prompt = f"{user_prompt}\n\n{fallback_prompt}".strip()

    # 2. Build Gemini content parts from either the uploaded image or extracted text.
    content_parts: list[types.Part] = []

    if content_type.startswith("image/"):
        content_parts = [
            types.Part.from_text(text=final_prompt),
            types.Part.from_bytes(data=file_bytes, mime_type=content_type),
        ]

    elif content_type == "application/pdf" or filename.endswith(".pdf"):
        extracted_text = extract_text_from_pdf(file_bytes)
        content_parts = [types.Part.from_text(text=f"{final_prompt}\n\n=== SOURCE DOCUMENT ===\n{extracted_text}")]

    elif filename.endswith(".docx"):
        extracted_text = extract_text_from_docx(file_bytes)
        content_parts = [types.Part.from_text(text=f"{final_prompt}\n\n=== SOURCE DOCUMENT ===\n{extracted_text}")]

    elif filename.endswith(".pptx"):
        extracted_text = extract_text_from_pptx(file_bytes)
        content_parts = [types.Part.from_text(text=f"{final_prompt}\n\n=== SOURCE DOCUMENT ===\n{extracted_text}")]

    elif content_type.startswith("text/") or filename.endswith((".txt", ".md", ".json")):
        try:
            extracted_text = file_bytes.decode("utf-8")
            content_parts = [types.Part.from_text(text=f"{final_prompt}\n\n=== SOURCE DOCUMENT ===\n{extracted_text}")]
        except UnicodeDecodeError:
            raise HTTPException(status_code=400, detail="Unable to decode text file as UTF-8.")
    else:
        raise HTTPException(status_code=400, detail="Unsupported file format.")

    try:
        response = await get_gemini_client().aio.models.generate_content(
            model=MODEL_NAME,
            contents=[types.Content(role="user", parts=content_parts)],
            config=types.GenerateContentConfig(
                system_instruction=[types.Part.from_text(text=system_instruction)],
                thinking_config=types.ThinkingConfig(thinking_level="HIGH"),
            ),
        )
        if not response.text:
            raise HTTPException(status_code=502, detail="Gemma returned an empty explanation.")
        return {"text": response.text.strip()}
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=502, detail=f"Gemma explanation generation failed: {exc}")

@app.post("/quiz")
async def generate_quiz(payload: QuizRequest):
    """Generate a quiz with the Google AI Studio Gemma 4 configuration."""
    system_instruction = """You are an expert assessment generator. Analyze the academic study guide and generate 3 highly relevant, concept-testing multiple-choice questions.

Output ONLY valid JSON in this exact shape:
{
  "questions": [
    {
      "id": 1,
      "question": "Clear, contextual question in the same language as the study guide.",
      "options": ["Option A", "Option B", "Option C", "Option D"],
      "answer": "The exact correct option text",
      "explanation": "A short explanation of why the answer is correct."
    }
  ]
}

Do not use Markdown fences or any text outside the JSON object."""

    try:
        response = await get_gemini_client().aio.models.generate_content(
            model=MODEL_NAME,
            contents=[types.Content(role="user", parts=[
                types.Part.from_text(text=f"Study guide:\n{payload.text}"),
            ])],
            config=types.GenerateContentConfig(
                thinking_config=types.ThinkingConfig(thinking_level="HIGH"),
                tools=[types.Tool(google_search=types.GoogleSearch())],
                system_instruction=[types.Part.from_text(text=system_instruction)],
            ),
        )
        generated = json.loads((response.text or "").strip())
        questions = generated.get("questions", [])
        if not isinstance(questions, list) or not questions:
            raise ValueError("The response did not include a questions array.")

        # Preserve the existing frontend contract.
        return {"quiz": [
            {
                "question": question["question"],
                "options": question["options"],
                "correct_answer": question["answer"],
            }
            for question in questions
        ]}
    except json.JSONDecodeError as exc:
        raise HTTPException(status_code=502, detail=f"Gemma returned invalid quiz JSON: {exc}")
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=502, detail=f"Gemma quiz generation failed: {exc}")


@app.post("/study_plan")
async def generate_study_plan(payload: StudyPlanRequest):
    """Generate a study plan in the response shape expected by the dashboard."""
    context = "\n\n".join(
        str(message.get("content", ""))
        for message in payload.messages
        if isinstance(message, dict) and message.get("content")
    )
    if not context:
        raise HTTPException(status_code=400, detail="Study-plan context is required.")

    system_instruction = """You are an educational planning assistant. Create a practical, concise study plan and exactly three actionable reminders from the supplied course context.

Return ONLY valid JSON in this exact shape:
{
  "planMarkdown": "A clear step-by-step study plan.",
  "reminders": ["Reminder one", "Reminder two", "Reminder three"]
}
Do not use Markdown fences or add text outside the JSON object."""

    try:
        response = await get_gemini_client().aio.models.generate_content(
            model=MODEL_NAME,
            contents=[types.Content(role="user", parts=[types.Part.from_text(text=context)])],
            config=types.GenerateContentConfig(
                thinking_config=types.ThinkingConfig(thinking_level="HIGH"),
                system_instruction=[types.Part.from_text(text=system_instruction)],
            ),
        )
        plan = json.loads((response.text or "").strip())
        if not isinstance(plan.get("planMarkdown"), str) or not isinstance(plan.get("reminders"), list):
            raise ValueError("The response did not have the expected study-plan structure.")

        # Retain the existing dashboard's `data.message.content` contract.
        return {"message": {"content": json.dumps(plan)}}
    except json.JSONDecodeError as exc:
        raise HTTPException(status_code=502, detail=f"Gemma returned invalid study-plan JSON: {exc}")
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=502, detail=f"Gemma study-plan generation failed: {exc}")
